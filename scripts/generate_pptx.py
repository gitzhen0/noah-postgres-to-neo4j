"""
Generate a submission PPTX mirroring the HTML deck at presentation/index.html.

Live-demo slides (7/9/11/14) are replaced with frozen artifacts: captured
terminal output, proof tables, and screenshot placeholders. The spoken
presentation will still use the HTML deck; this file is the "one printable
artifact" required for submission.

Run:
    ./venv/bin/python scripts/generate_pptx.py
Outputs to:
    ./NOAH_Capstone_Final_Presentation.pptx
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── theme ──────────────────────────────────────────────────────────────────
INK        = RGBColor(0x0F, 0x17, 0x2A)   # near-black
INK2       = RGBColor(0x33, 0x41, 0x55)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
RULE       = RGBColor(0xE2, 0xE8, 0xF0)
BG         = RGBColor(0xFF, 0xFF, 0xFF)
BG_SUBTLE  = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)   # blue
ACCENT_BG  = RGBColor(0xEF, 0xF6, 0xFF)
GREEN      = RGBColor(0x05, 0x96, 0x69)
GREEN_BG   = RGBColor(0xEC, 0xFD, 0xF5)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
AMBER_BG   = RGBColor(0xFF, 0xFB, 0xEB)
RED        = RGBColor(0xDC, 0x26, 0x26)
TERM_BG    = RGBColor(0x0F, 0x17, 0x2A)
TERM_FG    = RGBColor(0xE2, 0xE8, 0xF0)
TERM_MUTED = RGBColor(0x94, 0xA3, 0xB8)
TERM_GREEN = RGBColor(0x6E, 0xE7, 0xB7)
TERM_AMBER = RGBColor(0xFD, 0xE6, 0x8A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.55)

FONT_BODY = "Calibri"
FONT_MONO = "Consolas"


# ── helpers ────────────────────────────────────────────────────────────────
def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_run(run, *, size=None, bold=None, italic=None, color=None, font=None):
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if font is not None:
        run.font.name = font


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, italic=italic, color=color, font=font)
    return tb


def add_rect(slide, left, top, width, height, *, fill=None, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    # Smaller corner radius for rounded rects
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, AttributeError):
            pass
    return shape


def add_eyebrow(slide, text, *, top=MARGIN):
    add_text(slide, MARGIN, top, Inches(12), Inches(0.3), text.upper(),
             size=11, bold=True, color=MUTED)


def add_title(slide, text, *, top=Inches(0.95), accent_word=None):
    """Add a section-style heading. If accent_word given, color that part blue."""
    tb = slide.shapes.add_textbox(MARGIN, top, Inches(12), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if accent_word and accent_word in text:
        before, rest = text.split(accent_word, 1)
        for segment, color, italic in [(before, INK, False),
                                        (accent_word, ACCENT, False),
                                        (rest, INK, False)]:
            if not segment:
                continue
            r = p.add_run()
            r.text = segment
            set_run(r, size=32, bold=True, italic=italic, color=color, font=FONT_BODY)
    else:
        r = p.add_run()
        r.text = text
        set_run(r, size=32, bold=True, color=INK, font=FONT_BODY)
    return tb


def add_footer(slide, page_num, total):
    add_text(slide, MARGIN, Inches(7.05), Inches(6), Inches(0.3),
             "NOAH RDBMS → Knowledge Graph · Zhen Yang · NYU SPS MASY Capstone 2026",
             size=9, color=MUTED)
    add_text(slide, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.3),
             f"{page_num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=INK2, spacing=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.level = 0
        r = p.add_run()
        r.text = f"•  {item}"
        set_run(r, size=size, color=color, font=FONT_BODY)


def add_terminal(slide, left, top, width, height, lines, *, font_size=11):
    """Render a code/terminal block. `lines` is list of (text, color) tuples."""
    add_rect(slide, left, top, width, height, fill=TERM_BG, radius=True)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12),
                                   width - Inches(0.3), height - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    for i, (text, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = text
        set_run(r, size=font_size, color=color, font=FONT_MONO)


def add_table(slide, left, top, width, height, rows, cols, *,
              header=True, col_widths=None):
    gt = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = w
    return gt


def set_cell(cell, text, *, size=11, bold=False, color=INK, fill=None,
             align=PP_ALIGN.LEFT, font=FONT_BODY):
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color, font=font)


# ── slides ─────────────────────────────────────────────────────────────────

def slide_01_title(prs, n, total):
    s = add_blank_slide(prs)
    # accent strip on top
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill=ACCENT)
    add_text(s, MARGIN, Inches(1.2), Inches(12), Inches(0.4),
             "NYU SPS MASY · Capstone · Spring 2026".upper(),
             size=12, bold=True, color=ACCENT)
    add_text(s, MARGIN, Inches(1.8), Inches(12), Inches(2.0),
             "Automated RDBMS → Knowledge Graph Conversion",
             size=46, bold=True, color=INK)
    add_text(s, MARGIN, Inches(3.2), Inches(12), Inches(1.6),
             "A natural-language interface over NYC affordable housing data, built on "
             "8,604 real housing projects and validated end-to-end against PostgreSQL + PostGIS.",
             size=17, color=MUTED)
    # signature block
    add_rect(s, MARGIN, Inches(5.6), Inches(0.04), Inches(1.2), fill=ACCENT)
    add_text(s, MARGIN + Inches(0.25), Inches(5.6), Inches(6), Inches(0.4),
             "Zhen Yang", size=16, bold=True, color=INK)
    add_text(s, MARGIN + Inches(0.25), Inches(6.0), Inches(6), Inches(0.4),
             "Project Sponsor: Dr. Andres Fortino · Digital Forge Lab",
             size=12, color=MUTED)
    add_text(s, MARGIN + Inches(0.25), Inches(6.35), Inches(6), Inches(0.4),
             "Final Defense · April 2026", size=12, color=MUTED)
    add_footer(s, n, total)


def slide_02_problem(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "The Problem")
    add_title(s, "Relational databases optimize for storage, not for relationships.",
              accent_word="storage")

    # Left: SQL
    sql_top = Inches(2.4)
    add_text(s, MARGIN, Inches(2.05), Inches(6.3), Inches(0.3),
             "PostgreSQL  ·  21 lines  ·  174.4 ms",
             size=11, bold=True, color=MUTED)
    add_terminal(s, MARGIN, sql_top, Inches(6.3), Inches(3.8), [
        ("-- \"All ZIPs within 2 hops of 10001\"", TERM_MUTED),
        ("WITH RECURSIVE reachable(zip_code, depth) AS (", TERM_FG),
        ("  SELECT zip_code, 0", TERM_FG),
        ("  FROM   zip_shapes", TERM_FG),
        ("  WHERE  zip_code = '10001'", TERM_FG),
        ("  UNION", TERM_FG),
        ("  SELECT CASE WHEN a.zip_code = r.zip_code", TERM_FG),
        ("              THEN b.zip_code ELSE a.zip_code END,", TERM_FG),
        ("         r.depth + 1", TERM_FG),
        ("  FROM   reachable r", TERM_FG),
        ("  JOIN   zip_shapes a JOIN zip_shapes b", TERM_FG),
        ("         ON  a.zip_code < b.zip_code", TERM_FG),
        ("         AND ST_Touches(a.geom, b.geom)", TERM_FG),
        ("    ON   r.zip_code IN (a.zip_code, b.zip_code)", TERM_FG),
        ("  WHERE  r.depth < 2", TERM_FG),
        (")", TERM_FG),
        ("SELECT DISTINCT zip_code, MIN(depth)", TERM_FG),
        ("FROM reachable GROUP BY zip_code;", TERM_FG),
    ], font_size=10)

    # Right: Cypher
    add_text(s, Inches(7.0), Inches(2.05), Inches(6.0), Inches(0.3),
             "Neo4j  ·  4 lines  ·  4.6 ms",
             size=11, bold=True, color=ACCENT)
    add_terminal(s, Inches(7.0), sql_top, Inches(5.8), Inches(3.8), [
        ("// same question", TERM_MUTED),
        ("MATCH path = (z:ZipCode {zip_code:'10001'})", TERM_FG),
        ("              -[:NEIGHBORS*0..2]-(n:ZipCode)", TERM_FG),
        ("WITH  n.zip_code AS zip_code,", TERM_FG),
        ("      min(length(path)) AS min_depth", TERM_FG),
        ("RETURN zip_code, min_depth", TERM_FG),
        ("ORDER BY min_depth, zip_code;", TERM_FG),
    ], font_size=11)

    add_text(s, MARGIN, Inches(6.5), Inches(12), Inches(0.5),
             "Same question. Different paradigm.  5× less code  ·  37× faster runtime.",
             size=14, bold=True, color=INK)
    add_footer(s, n, total)


def slide_03_goal(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Goal")
    add_title(s, "Build a reusable conversion bot — and prove it on real data.",
              accent_word="reusable")
    add_text(s, MARGIN, Inches(2.1), Inches(12), Inches(1.2),
             "Convert the NOAH (Naturally Occurring Affordable Housing) PostgreSQL "
             "database — 8,604 projects, 177 NYC ZIPs, 2,225 census tracts, PostGIS "
             "geometry — into a Neo4j knowledge graph, with a natural-language query "
             "interface that non-technical users can actually use.",
             size=15, color=INK2)

    # 4 stat cards
    cards = [
        ("8,604",  "housing projects migrated"),
        ("≥ 75 %", "Text2Cypher accuracy target"),
        ("0",      "data loss on referential integrity"),
        ("12 weeks","from proposal to final delivery"),
    ]
    card_w = Inches(2.8)
    gap = Inches(0.25)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(4.1)
    for i, (num, label) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(1.7), fill=BG_SUBTLE, line=RULE, radius=True)
        add_text(s, x, top + Inches(0.2), card_w, Inches(0.8),
                 num, size=32, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), top + Inches(1.05), card_w - Inches(0.3),
                 Inches(0.6), label, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(s, MARGIN, Inches(6.3), Inches(12), Inches(0.4),
             "Four SMART objectives from the project specification, each with a "
             "concrete validation artifact committed in the repo.",
             size=12, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_04_papers(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Built on published research")
    add_title(s, "Three papers, synthesized into a production-ready pipeline.")

    papers = [
        ("FORMAL MODEL",
         "De Virgilio et al., 2013",
         "ACM GRADES Workshop",
         "Formal conversion rules: tables → nodes, foreign keys → relationships, "
         "join tables → direct edges. Exploits RDBMS constraints to preserve "
         "semantic integrity."),
        ("AUTOMATION",
         "Rel2Graph",
         "Zhao et al., 2023  ·  arXiv:2310.01080",
         "Automated knowledge graph construction from multiple relational "
         "sources, validated on Spider / KaggleDBQA benchmarks. Proves the "
         "pipeline shape scales beyond a single schema."),
        ("TOOLING",
         "Data2Neo",
         "Minder et al., 2024  ·  arXiv:2406.04995",
         "Open-source Python patterns for incremental Neo4j loading with "
         "configurable mapping. Operational blueprint for batch MERGE with "
         "idempotency."),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(2.3)
    for i, (tag, title, authors, body) in enumerate(papers):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(3.9), fill=BG_SUBTLE, line=RULE, radius=True)
        add_text(s, x + Inches(0.25), top + Inches(0.2), card_w - Inches(0.4), Inches(0.3),
                 tag, size=9, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), top + Inches(0.55), card_w - Inches(0.4), Inches(0.5),
                 title, size=16, bold=True, color=INK)
        add_text(s, x + Inches(0.25), top + Inches(1.0), card_w - Inches(0.4), Inches(0.4),
                 authors, size=10, color=MUTED)
        add_text(s, x + Inches(0.25), top + Inches(1.5), card_w - Inches(0.4), Inches(2.3),
                 body, size=11, color=INK2)

    add_text(s, MARGIN, Inches(6.4), Inches(12), Inches(0.5),
             "My contribution: synthesize all three into a single config-driven pipeline, "
             "add an LLM-powered schema interpreter, and validate on a real PostGIS dataset.",
             size=12, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_05_architecture(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "System Architecture")
    add_title(s, "Five components. One config file.", accent_word="One")

    # pipeline boxes left-to-right
    stages = [
        ("PostgreSQL",        "source",        RGBColor(0x33, 0x67, 0x91)),
        ("Schema Analyzer",   "metadata →",    ACCENT),
        ("LLM Interpreter",   "→ draft YAML",  AMBER),
        ("Mapping Engine",    "YAML rules",    ACCENT),
        ("Data Migrator",     "MERGE loader",  ACCENT),
        ("Neo4j + Auditor",   "target",        GREEN),
    ]
    box_w = Inches(1.85)
    gap = Inches(0.18)
    total_w = box_w * 6 + gap * 5
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(2.8)
    for i, (title, sub, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, top, box_w, Inches(1.6), fill=color, radius=True)
        add_text(s, x + Inches(0.1), top + Inches(0.5), box_w - Inches(0.2), Inches(0.4),
                 title, size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), top + Inches(0.95), box_w - Inches(0.2), Inches(0.3),
                 sub, size=10, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        # arrow between
        if i < 5:
            add_text(s, x + box_w, top + Inches(0.55), gap, Inches(0.5),
                     "→", size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # key points below
    add_bullets(s, MARGIN, Inches(5.2), Inches(12), Inches(1.6), [
        "Everything except the Data Migrator is stateless.",
        "Schema Analyzer introspects PG metadata (tables, FKs, PostGIS types).",
        "LLM Interpreter produces a draft mapping_rules.yaml that a human reviews.",
        "None of the components contain NOAH-specific logic — the YAML is the only input that changes across datasets.",
    ], size=13)
    add_footer(s, n, total)


def slide_06_schema(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "The Target Model")
    add_title(s, "5 node labels, 6 relationship types, PostGIS-native.",
              accent_word="PostGIS-native")

    # nodes table
    tbl = add_table(s, MARGIN, Inches(2.2), Inches(5.9), Inches(3.5), 6, 3,
                    col_widths=[Inches(2.4), Inches(1.1), Inches(2.4)])
    headers = ["Node label", "Count", "Notes"]
    for i, h in enumerate(headers):
        set_cell(tbl.cell(0, i), h, size=11, bold=True, color=INK, fill=BG_SUBTLE)
    rows = [
        ("HousingProject",        "8,604", "Socrata hg8x-zxpr"),
        ("ZipCode",                 "177", "PostGIS geometry"),
        ("Demographic",             "176", "ACS 2022"),
        ("AffordabilityAnalysis",   "177", "per-ZIP rollup"),
        ("RentBurden",            "2,225", "census tract-level"),
    ]
    for r, (lbl, cnt, note) in enumerate(rows, start=1):
        set_cell(tbl.cell(r, 0), lbl, size=11, bold=True, color=INK)
        set_cell(tbl.cell(r, 1), cnt, size=11, color=INK2, align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(r, 2), note, size=10, color=MUTED)

    # relationships table
    tbl2 = add_table(s, Inches(7.0), Inches(2.2), Inches(5.9), Inches(3.9), 7, 3,
                     col_widths=[Inches(2.4), Inches(1.1), Inches(2.4)])
    for i, h in enumerate(["Relationship type", "Count", "Source"]):
        set_cell(tbl2.cell(0, i), h, size=11, bold=True, color=INK, fill=BG_SUBTLE)
    rel_rows = [
        ("LOCATED_IN_ZIP",         "6,851", "FK (postcode)"),
        ("HAS_DEMOGRAPHICS",         "176", "FK · new in v2"),
        ("HAS_AFFORDABILITY_DATA",   "177", "FK"),
        ("IN_CENSUS_TRACT",        "5,426", "computed"),
        ("NEIGHBORS",                "392", "spatial · bidir."),
        ("CONTAINS_TRACT",         "4,050", "spatial"),
    ]
    for r, (t, c, src) in enumerate(rel_rows, start=1):
        set_cell(tbl2.cell(r, 0), t, size=11, bold=True, color=INK, font=FONT_MONO)
        set_cell(tbl2.cell(r, 1), c, size=11, color=INK2, align=PP_ALIGN.RIGHT)
        set_cell(tbl2.cell(r, 2), src, size=10, color=MUTED)

    add_text(s, MARGIN, Inches(6.4), Inches(12), Inches(0.5),
             "ZipCode sits at the hub. NEIGHBORS is self-referential and bidirectional, "
             "enabling multi-hop spatial queries — the central win of the graph model.",
             size=12, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_07_live_migration(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Live Demo · 1 of 4")
    add_title(s, "Clean Neo4j → full graph in ~10 seconds.", accent_word="~10 seconds")

    # live note bar
    add_rect(s, MARGIN, Inches(2.1), Inches(12.2), Inches(0.5),
             fill=AMBER_BG, line=RGBColor(0xFD, 0xE6, 0x8A), radius=True)
    add_text(s, MARGIN + Inches(0.25), Inches(2.16), Inches(11.7), Inches(0.4),
             "During the live defense: Neo4j Browser shows count = 0  →  terminal runs migrate  →  "
             "Neo4j Browser shows count = 11,183 and renders the graph.",
             size=11, bold=True, color=AMBER)

    # left: migrate output
    add_text(s, MARGIN, Inches(2.9), Inches(6.3), Inches(0.3),
             "STEP 1   Migrate PG → Neo4j", size=11, bold=True, color=INK)
    add_terminal(s, MARGIN, Inches(3.25), Inches(6.3), Inches(3.3), [
        ("$ python main.py migrate --clear \\", TERM_GREEN),
        ("    --mapping-rules outputs/mapping_draft.yaml", TERM_GREEN),
        ("[INFO] Clearing Neo4j database...", TERM_MUTED),
        ("[INFO] Setting up constraints + indexes", TERM_MUTED),
        ("✓ HousingProject         8,604 nodes", TERM_FG),
        ("✓ ZipCode                  177  ·  Demographic 176", TERM_FG),
        ("✓ AffordabilityAnalysis    177  ·  RentBurden 2,225", TERM_FG),
        ("✓ LOCATED_IN_ZIP         6,851  (35 orphan skipped)", TERM_FG),
        ("✓ HAS_DEMOGRAPHICS         176  ← new in v2", TERM_FG),
        ("✓ HAS_AFFORDABILITY_DATA   177", TERM_FG),
        ("✓ IN_CENSUS_TRACT        5,426", TERM_FG),
        ("✓ NEIGHBORS                392  ·  CONTAINS_TRACT 4,050", TERM_FG),
        ("", TERM_FG),
        ("11,359 nodes · 17,072 relationships", TERM_AMBER),
        ("real  7.89 s", TERM_MUTED),
    ], font_size=10)

    # right: audit output
    add_text(s, Inches(7.0), Inches(2.9), Inches(5.8), Inches(0.3),
             "STEP 2   Verify integrity", size=11, bold=True, color=INK)
    add_terminal(s, Inches(7.0), Inches(3.25), Inches(5.8), Inches(3.3), [
        ("$ python main.py audit", TERM_GREEN),
        ("Node Count Verification", TERM_MUTED),
        ("✓ MATCH   all 5 labels, PG↔Neo4j parity", TERM_FG),
        ("", TERM_FG),
        ("Relationship Counts", TERM_MUTED),
        ("✓ MATCH   LOCATED_IN_ZIP    6,851 / 6,851", TERM_FG),
        ("✓ MATCH   HAS_DEMOGRAPHICS    176 /   176", TERM_FG),
        ("✓ MATCH   HAS_AFFORDABILITY   177 /   177", TERM_FG),
        ("✓ MATCH   + 3 computed/spatial rels", TERM_FG),
        ("", TERM_FG),
        ("Data Sample Verification (20 rows per label)", TERM_MUTED),
        ("✓ 100.0%  on all 5 labels", TERM_FG),
        ("", TERM_FG),
        ("INFO   LOCATED_IN_ZIP — 35 orphan FKs", TERM_AMBER),
        ("       (expected; postcodes outside NYC)", TERM_MUTED),
        ("Overall Status: PASS", TERM_GREEN),
    ], font_size=10)

    add_text(s, MARGIN, Inches(6.75), Inches(12), Inches(0.4),
             "Total: migrate 7.89 s + audit 1.32 s ≈ ~10 seconds for 11k nodes / 17k "
             "relationships / 100 % integrity / sample-verified. Idempotent — safe to rerun.",
             size=11, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_08_validation(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Validation")
    add_title(s, "Every relationship accounted for.")

    tbl = add_table(s, MARGIN, Inches(2.2), Inches(12.2), Inches(3.5), 7, 6,
                    col_widths=[Inches(3.0), Inches(1.4), Inches(2.2),
                                Inches(1.8), Inches(1.8), Inches(2.0)])
    for i, h in enumerate(["Relationship", "Source", "PostgreSQL",
                            "Neo4j", "Orphans", "Status"]):
        align = PP_ALIGN.RIGHT if i in (2, 3, 4) else PP_ALIGN.LEFT
        set_cell(tbl.cell(0, i), h, size=11, bold=True, color=INK,
                 fill=BG_SUBTLE, align=align)

    rows = [
        ("LOCATED_IN_ZIP",           "FK",        "6,851 expected", "6,851", "35 (INFO)", "✓ MATCH"),
        ("HAS_DEMOGRAPHICS",          "FK · new",   "176",           "176",   "0",          "✓ MATCH"),
        ("HAS_AFFORDABILITY_DATA",   "FK",        "177",           "177",   "0",          "✓ MATCH"),
        ("IN_CENSUS_TRACT",          "Computed",   "N/A",           "5,426", "—",          "✓ OK"),
        ("NEIGHBORS",                 "Spatial",   "N/A",           "392",   "—",          "✓ OK"),
        ("CONTAINS_TRACT",           "Spatial",   "N/A",           "4,050", "—",          "✓ OK"),
    ]
    for r, (name, src, pg, n4j, orp, status) in enumerate(rows, start=1):
        set_cell(tbl.cell(r, 0), name, size=11, bold=True, color=INK, font=FONT_MONO)
        set_cell(tbl.cell(r, 1), src, size=10, color=INK2)
        set_cell(tbl.cell(r, 2), pg,  size=11, color=INK2, align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(r, 3), n4j, size=11, color=INK,  align=PP_ALIGN.RIGHT, bold=True)
        set_cell(tbl.cell(r, 4), orp, size=11, color=INK2, align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(r, 5), status, size=11, bold=True, color=GREEN, fill=GREEN_BG)

    add_text(s, MARGIN, Inches(5.9), Inches(12.2), Inches(1.2),
             "On the \"35 orphans\": these are housing projects whose postcode falls outside "
             "NYC's 177-ZIP coverage (adjacent NJ/PA, PO boxes). MERGE correctly skips them. "
             "Surfaced as INFO, not WARN — that distinction is the difference between an audit "
             "that cries wolf and an audit that tells the truth.",
             size=12, color=INK2)
    add_footer(s, n, total)


def slide_09_text2cypher(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Live Demo · 2 of 4")
    add_title(s, "Natural language → Cypher → graph.")

    add_text(s, MARGIN, Inches(2.1), Inches(12), Inches(0.6),
             "The Streamlit \"Ask\" page takes English, generates Cypher via Claude, executes "
             "it against Neo4j, and explains results in plain English — all in one round-trip.",
             size=13, color=INK2)

    # Example block
    add_rect(s, MARGIN, Inches(3.0), Inches(12.2), Inches(3.5),
             fill=BG_SUBTLE, line=RULE, radius=True)
    add_text(s, MARGIN + Inches(0.3), Inches(3.15), Inches(11.5), Inches(0.35),
             "User input", size=10, bold=True, color=MUTED)
    add_text(s, MARGIN + Inches(0.3), Inches(3.5), Inches(11.5), Inches(0.5),
             "\"How many housing projects are in each borough?\"",
             size=16, italic=True, color=INK)

    add_text(s, MARGIN + Inches(0.3), Inches(4.15), Inches(11.5), Inches(0.35),
             "Generated Cypher  (via claude-sonnet-4-5)", size=10, bold=True, color=MUTED)
    add_terminal(s, MARGIN + Inches(0.3), Inches(4.5), Inches(11.5), Inches(1.0), [
        ("MATCH (h:HousingProject) WHERE h.borough IS NOT NULL", TERM_FG),
        ("RETURN h.borough AS borough, count(h) AS projects", TERM_FG),
        ("ORDER BY projects DESC;", TERM_FG),
    ], font_size=12)
    add_text(s, MARGIN + Inches(0.3), Inches(5.65), Inches(11.5), Inches(0.75),
             "Brooklyn 2,414  ·  Bronx 2,093  ·  Manhattan 1,833  ·  "
             "Queens 1,612  ·  Staten Island 652",
             size=13, color=INK)

    add_text(s, MARGIN, Inches(6.6), Inches(12), Inches(0.5),
             "Live: http://localhost:8505/Ask  —  embedded as iframe inside the HTML deck.",
             size=11, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_10_accuracy(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Benchmarking the NL Interface")
    add_title(s, "Text2Cypher: 95% on a 20-question benchmark.",
              accent_word="95%")

    # Big number card
    add_rect(s, MARGIN, Inches(2.3), Inches(5.5), Inches(3.5),
             fill=GREEN_BG, line=GREEN, radius=True)
    add_text(s, MARGIN, Inches(2.5), Inches(5.5), Inches(1.5),
             "95 %", size=110, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(4.5), Inches(5.5), Inches(0.5),
             "19 of 20 correct", size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.0), Inches(5.5), Inches(0.6),
             "(75% target · 20 ppts over plan)",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # Breakdown by difficulty
    tbl = add_table(s, Inches(6.5), Inches(2.3), Inches(6.3), Inches(3.2), 4, 3,
                    col_widths=[Inches(2.3), Inches(1.8), Inches(2.2)])
    for i, h in enumerate(["Difficulty", "Questions", "Pass rate"]):
        set_cell(tbl.cell(0, i), h, size=11, bold=True, color=INK, fill=BG_SUBTLE)
    difficulty_rows = [
        ("Easy",    "6",  "6 / 6  = 100 %"),
        ("Medium",  "10", "10 / 10 = 100 %"),
        ("Hard",    "4",  "3 / 4  =  75 %"),
    ]
    for r, (d, q, p) in enumerate(difficulty_rows, start=1):
        set_cell(tbl.cell(r, 0), d, size=12, bold=True, color=INK)
        set_cell(tbl.cell(r, 1), q, size=12, color=INK2, align=PP_ALIGN.CENTER)
        set_cell(tbl.cell(r, 2), p, size=12, color=INK2, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(6.5), Inches(5.7), Inches(6.3), Inches(1.3),
             "The only miss was a hard-difficulty question requiring multi-clause subqueries. "
             "Grading: 4 dimensions per question — schema correctness, syntactic validity, "
             "execution result parity vs ground truth, and explanation readability.",
             size=11, color=INK2)
    add_footer(s, n, total)


def slide_11_graphviz(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Live Demo · 3 of 4")
    add_title(s, "From query to picture — variable-length paths, visualized.")

    add_text(s, MARGIN, Inches(2.1), Inches(12), Inches(0.6),
             "The Streamlit \"Explore\" page executes ad-hoc Cypher and renders results as "
             "an interactive force-directed graph (pyvis). Useful for ZIP-to-ZIP neighborhood "
             "traversals that are hard to intuit from tabular output.",
             size=13, color=INK2)

    add_text(s, MARGIN, Inches(3.1), Inches(12), Inches(0.35),
             "Example query", size=11, bold=True, color=MUTED)
    add_terminal(s, MARGIN, Inches(3.5), Inches(12.2), Inches(1.4), [
        ("MATCH path = (z:ZipCode {zip_code:'10001'})", TERM_FG),
        ("              -[:NEIGHBORS*0..2]-(n:ZipCode)", TERM_FG),
        ("RETURN path LIMIT 50;", TERM_FG),
    ], font_size=14)

    # Visualization placeholder
    add_rect(s, MARGIN, Inches(5.15), Inches(12.2), Inches(1.8),
             fill=ACCENT_BG, line=ACCENT, radius=True)
    add_text(s, MARGIN + Inches(0.35), Inches(5.3), Inches(11.5), Inches(0.45),
             "Rendered inside Streamlit Explore  ·  http://localhost:8505/Explore",
             size=11, bold=True, color=ACCENT)
    add_bullets(s, MARGIN + Inches(0.35), Inches(5.75), Inches(11.5), Inches(1.2), [
        "Starts at ZIP 10001 (Manhattan), expands outward by shared-boundary edges.",
        "Each hop recomputed in <5 ms; visualization refreshes in <1 s for 50-node path.",
        "In SQL this is a 21-line WITH RECURSIVE; in Cypher it is 4 lines.",
    ], size=11)
    add_footer(s, n, total)


def slide_12_perf(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Performance  ·  PostgreSQL vs Neo4j  ·  10 queries")
    add_title(s, "Right tool for the right query shape.")

    tbl = add_table(s, MARGIN, Inches(2.2), Inches(12.2), Inches(4.0), 6, 5,
                    col_widths=[Inches(3.5), Inches(1.5), Inches(2.2),
                                Inches(2.2), Inches(2.8)])
    for i, h in enumerate(["Query category", "Count", "PG median",
                            "Neo4j median", "Verdict"]):
        align = PP_ALIGN.RIGHT if i in (2, 3) else PP_ALIGN.LEFT
        set_cell(tbl.cell(0, i), h, size=11, bold=True, color=INK,
                 fill=BG_SUBTLE, align=align)

    cats = [
        ("Aggregation (GROUP BY)",          "3", "8 ms",    "24 ms",    "PG wins (4.4×)"),
        ("Single-hop FK traversal",         "2", "14 ms",   "9 ms",     "Neo4j wins (1.6×)"),
        ("Multi-hop / recursive traversal", "2", "174 ms",  "5 ms",     "Neo4j wins (37×)"),
        ("Spatial ST_* predicate",          "2", "61 ms",   "18 ms",    "Neo4j wins (3.4×)"),
        ("Analytical with CTE",             "1", "112 ms",  "32 ms",    "Neo4j wins (3.5×)"),
    ]
    for r, (cat, cnt, pg, n4j, verdict) in enumerate(cats, start=1):
        set_cell(tbl.cell(r, 0), cat, size=11, color=INK)
        set_cell(tbl.cell(r, 1), cnt, size=11, color=INK2, align=PP_ALIGN.CENTER)
        set_cell(tbl.cell(r, 2), pg,  size=11, color=INK2, align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(r, 3), n4j, size=11, color=INK, align=PP_ALIGN.RIGHT, bold=True)
        win_color = MUTED if "PG wins" in verdict else GREEN
        set_cell(tbl.cell(r, 4), verdict, size=11, bold=True, color=win_color)

    add_text(s, MARGIN, Inches(6.4), Inches(12), Inches(0.5),
             "Neo4j is not universally faster. It wins on shapes SQL struggles with — "
             "variable-length traversal, pre-materialized spatial adjacency — and loses on "
             "bulk aggregation. The honest claim is \"right tool for query shape,\" not \"always faster.\"",
             size=12, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_13_hero(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "The query that defines the architecture")
    add_title(s, "\"All ZIPs within 2 hops of 10001\" — 37× faster.",
              accent_word="37× faster")

    # big number
    add_rect(s, MARGIN, Inches(2.2), Inches(5.2), Inches(4.0),
             fill=GREEN_BG, line=GREEN, radius=True)
    add_text(s, MARGIN, Inches(2.6), Inches(5.2), Inches(1.2),
             "37.7 ×", size=100, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(4.5), Inches(5.2), Inches(0.5),
             "PG 174.4 ms  →  Neo4j 4.6 ms",
             size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.0), Inches(5.2), Inches(0.5),
             "median of 10 warm runs", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.45), Inches(5.2), Inches(0.5),
             "21 lines of SQL  →  4 lines of Cypher",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(s, Inches(6.5), Inches(2.2), Inches(6.3), Inches(0.4),
             "Why graph wins here", size=12, bold=True, color=ACCENT)
    add_bullets(s, Inches(6.5), Inches(2.7), Inches(6.3), Inches(4.0), [
        "Path-shaped query: expressed natively as a pattern, not approximated by joins.",
        "Pre-computed edges: NEIGHBORS materialized at migrate time; query is O(edges), not O(n²).",
        "No plan blowup: depth is a parameter, not a schema change. 3-hop = *0..3.",
        "Readable by a non-engineer in the room.",
    ], size=13)
    add_text(s, MARGIN, Inches(6.8), Inches(12), Inches(0.4),
             "Full benchmark: outputs/performance_report.json.  10 queries · 5 categories · warm cache.",
             size=10, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_14_agnostic(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Live Demo · 4 of 4  ·  Not just NOAH")
    add_title(s, "Three more databases. Same engine. Under 4 seconds each.",
              accent_word="Same engine")

    # proof table
    tbl = add_table(s, MARGIN, Inches(2.1), Inches(12.2), Inches(2.7), 5, 7,
                    col_widths=[Inches(1.5), Inches(3.1), Inches(1.0),
                                Inches(1.6), Inches(1.6), Inches(1.7), Inches(1.7)])
    for i, h in enumerate(["Dataset", "Domain", "Tables", "Nodes",
                           "Edges", "Migrate", "Orphans"]):
        align = PP_ALIGN.RIGHT if i >= 2 else PP_ALIGN.LEFT
        set_cell(tbl.cell(0, i), h, size=11, bold=True, color=INK,
                 fill=BG_SUBTLE, align=align)

    agn_rows = [
        ("noah",      "NYC housing · PostGIS  (reference)",     "5",  "11,359", "17,072", "7.89 s", "35 INFO*", True),
        ("chinook",   "Music store (Apple-iTunes shape)",         "11", "6,892",  "24,529", "2.71 s", "0",        False),
        ("northwind", "Wholesale orders (classic B2B)",           "8",  "1,050",  "4,807",  "1.52 s", "0",        False),
        ("pagila",    "DVD rental (Postgres official)",           "12", "25,758", "72,024", "3.36 s", "0",        False),
    ]
    for r, (name, dom, t, nd, ed, mig, orp, is_ref) in enumerate(agn_rows, start=1):
        fill = BG_SUBTLE if is_ref else None
        set_cell(tbl.cell(r, 0), name, size=11, bold=True, color=INK, font=FONT_MONO, fill=fill)
        set_cell(tbl.cell(r, 1), dom, size=10, color=INK2, fill=fill)
        set_cell(tbl.cell(r, 2), t,   size=11, color=INK2, align=PP_ALIGN.RIGHT, fill=fill)
        set_cell(tbl.cell(r, 3), nd,  size=11, color=INK2, align=PP_ALIGN.RIGHT, fill=fill)
        set_cell(tbl.cell(r, 4), ed,  size=11, color=INK2, align=PP_ALIGN.RIGHT, fill=fill)
        mig_color = GREEN if not is_ref else INK2
        set_cell(tbl.cell(r, 5), mig, size=11, bold=(not is_ref), color=mig_color,
                 align=PP_ALIGN.RIGHT, fill=fill)
        set_cell(tbl.cell(r, 6), orp, size=11, color=INK2, align=PP_ALIGN.RIGHT, fill=fill)

    # 3 what-made-each-interesting cards
    card_w = Inches(3.9)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(5.15)
    what = [
        ("chinook",   "Self-join recursion",      "Employee.reports_to becomes a REPORTS_TO edge. Org-chart traversal in Cypher: *0..n."),
        ("northwind", "Composite-PK junction",    "order_details collapsed into (Order)-[CONTAINS {qty,price}]->(Product)."),
        ("pagila",    "Layered geography",        "Address → City → Country plus two M2M junctions (film_actor, film_category)."),
    ]
    for i, (tag, title, body) in enumerate(what):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(1.5), fill=BG_SUBTLE, line=RULE, radius=True)
        add_text(s, x + Inches(0.2), top + Inches(0.1), card_w - Inches(0.3), Inches(0.3),
                 tag.upper(), size=9, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.2), top + Inches(0.4), card_w - Inches(0.3), Inches(0.4),
                 title, size=12, bold=True, color=INK)
        add_text(s, x + Inches(0.2), top + Inches(0.8), card_w - Inches(0.3), Inches(0.75),
                 body, size=10, color=INK2)

    add_text(s, MARGIN, Inches(6.85), Inches(12), Inches(0.4),
             "* NOAH's 35 orphans are legit out-of-NYC postcodes. Every FK in the 3 new "
             "datasets landed. Same src/noah_converter/ code — three ~200-line YAML files "
             "were the only input.",
             size=10, italic=True, color=MUTED)
    add_footer(s, n, total)


def slide_15_business(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Why this matters beyond the capstone")
    add_title(s, "Three audiences, one tool.")

    audiences = [
        ("Policy analysts & researchers",
         "Ask questions in English, get cited answers",
         [
             "Skip the SQL class — focus on the housing question.",
             "Cypher result comes with a paragraph explanation and a GeoJSON export.",
             "Schema-aware prompting prevents hallucinated column names.",
         ]),
        ("City data teams",
         "Drop-in migration pipeline for any PostgreSQL database",
         [
             "YAML-configured; no per-dataset code to maintain.",
             "Idempotent MERGE + audit report make re-migration safe.",
             "Proved on 4 unrelated schemas (NOAH + Chinook + Northwind + Pagila).",
         ]),
        ("Educators",
         "Complete teaching artifact",
         [
             "Jupyter notebook with graded labs (RDBMS → graph intuition).",
             "Public GitHub with Docker Compose; clone-and-run in 5 min.",
             "Open-source; Digital Forge Lab can reuse verbatim.",
         ]),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.2)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(2.3)
    for i, (name, tagline, bullets) in enumerate(audiences):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(4.2), fill=BG_SUBTLE, line=RULE, radius=True)
        add_text(s, x + Inches(0.25), top + Inches(0.2), card_w - Inches(0.4), Inches(0.5),
                 name, size=14, bold=True, color=INK)
        add_text(s, x + Inches(0.25), top + Inches(0.75), card_w - Inches(0.4), Inches(0.5),
                 tagline, size=11, italic=True, color=ACCENT)
        add_bullets(s, x + Inches(0.25), top + Inches(1.3), card_w - Inches(0.4),
                    Inches(2.8), bullets, size=10, spacing=4)
    add_footer(s, n, total)


def slide_16_limits(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Honest inventory")
    add_title(s, "What this project doesn't claim to solve.",
              accent_word="doesn't")

    add_bullets(s, MARGIN, Inches(2.2), Inches(12), Inches(4.5), [
        "Not real-time CDC. Migration is batch, re-runnable; streaming updates are out of scope.",
        "Not a general text-to-SQL/Cypher product. Schema-aware prompting works on this graph shape; open-domain NL querying over arbitrary schemas is a research problem.",
        "Not horizontally scaled. Tested on single-node Neo4j 5.15. No Fabric/cluster work.",
        "StreetEasy rent substitution. The originally planned median rent from StreetEasy was unavailable; replaced by ACS 2022 rent-burden rates. Documented in the audit report.",
        "35 FK orphans in LOCATED_IN_ZIP are real data-quality issues in the source dump — not a bug in the pipeline.",
        "Text2Cypher costs money per query (Claude API). Not zero-marginal-cost unless self-hosted.",
    ], size=13)
    add_footer(s, n, total)


def slide_17_lessons(prs, n, total):
    s = add_blank_slide(prs)
    add_eyebrow(s, "Three things I'd tell next semester's cohort")
    add_title(s, "Lessons learned.")

    lessons = [
        ("Audit semantics are the project.",
         "Early audits reported 35 expected skips as WARN. That's cry-wolf behavior. "
         "Distinguishing INFO (\"expected behavior\") from WARN (\"real data loss\") was a "
         "late-stage change that made the final report trustworthy. If the audit "
         "cannot be trusted, the pipeline cannot be trusted."),
        ("LLM = draftsperson, not decider.",
         "Claude produced the initial mapping YAML in 30 seconds. I spent an hour "
         "reviewing and tuning it. That ratio is correct. The LLM is a force-multiplier "
         "for the boilerplate; humans still own the semantic decisions."),
        ("Prove generalization early, not last.",
         "The \"works on Chinook / Northwind / Pagila\" demo was added in the final two "
         "weeks. If I had run that test in week 3, I would have caught three config "
         "assumptions that were silently NOAH-specific. Generalization is a test, not a claim."),
    ]
    top = Inches(2.2)
    for i, (head, body) in enumerate(lessons):
        y = top + Inches(i * 1.5)
        # number circle
        add_rect(s, MARGIN, y + Inches(0.08), Inches(0.5), Inches(0.5),
                 fill=ACCENT, radius=True)
        add_text(s, MARGIN, y + Inches(0.15), Inches(0.5), Inches(0.4),
                 str(i + 1), size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(s, MARGIN + Inches(0.8), y, Inches(11.4), Inches(0.4),
                 head, size=15, bold=True, color=INK)
        add_text(s, MARGIN + Inches(0.8), y + Inches(0.45), Inches(11.4), Inches(1.0),
                 body, size=11, color=INK2)
    add_footer(s, n, total)


def slide_18_qna(prs, n, total):
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill=ACCENT)

    add_text(s, MARGIN, Inches(1.2), Inches(12), Inches(0.4),
             "Thank you".upper(),
             size=12, bold=True, color=ACCENT)
    add_text(s, MARGIN, Inches(1.8), Inches(12), Inches(1.5),
             "Questions?", size=68, bold=True, color=INK)

    add_text(s, MARGIN, Inches(3.6), Inches(12), Inches(0.4),
             "Repository", size=12, bold=True, color=MUTED)
    add_text(s, MARGIN, Inches(4.0), Inches(12), Inches(0.5),
             "github.com/gitzhen0/noah-postgres-to-neo4j",
             size=15, color=ACCENT, font=FONT_MONO)

    add_text(s, MARGIN, Inches(4.9), Inches(12), Inches(0.4),
             "Evidence committed", size=12, bold=True, color=MUTED)
    add_bullets(s, MARGIN, Inches(5.3), Inches(12), Inches(2.0), [
        "outputs/audit_report.json  — machine-generated PASS record, timestamped",
        "outputs/benchmark_report.json  — 20-question Text2Cypher grading output",
        "outputs/performance_report.json  — 10-query benchmark, category_summary",
        "outputs/agnostic_benchmark.json  — 3-dataset generalization proof",
        "tests/unit/test_audit_semantics.py  — 9 unit tests covering audit logic",
        "docs/CAPSTONE_REPORT.pdf  — long-form write-up",
    ], size=12)
    add_footer(s, n, total)


# ── driver ─────────────────────────────────────────────────────────────────

def main():
    out = Path(__file__).resolve().parent.parent / "NOAH_Capstone_Final_Presentation.pptx"
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_01_title, slide_02_problem, slide_03_goal, slide_04_papers,
        slide_05_architecture, slide_06_schema, slide_07_live_migration,
        slide_08_validation, slide_09_text2cypher, slide_10_accuracy,
        slide_11_graphviz, slide_12_perf, slide_13_hero, slide_14_agnostic,
        slide_15_business, slide_16_limits, slide_17_lessons, slide_18_qna,
    ]
    total = len(slides)
    for i, fn in enumerate(slides, start=1):
        fn(prs, i, total)

    prs.save(str(out))
    print(f"Wrote: {out}")
    print(f"Size : {out.stat().st_size/1024:.1f} KB · {total} slides")


if __name__ == "__main__":
    main()
